# -*- coding: utf-8 -*-
"""Rebuild the Emotion Recognition presentation from REPORT.md (CNN-first)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = r"c:\Users\aleez\OneDrive\Desktop\5 sem uni\Emotion_Recognition_Project"
OUT = os.path.join(ROOT, "Emotion_Recognition_Presentation_CNN.pptx")
HIST = os.path.join(ROOT, "training", "training_history.png")
CONF = os.path.join(ROOT, "training", "confusion_matrix.png")

# ---------------------------------------------------------------- design system
NAVY      = RGBColor(0x26, 0x2B, 0x4A)
NAVY_DEEP = RGBColor(0x1E, 0x22, 0x3C)
CARD_D    = RGBColor(0x32, 0x38, 0x5C)
CARD_D2   = RGBColor(0x3A, 0x40, 0x68)
CREAM     = RGBColor(0xFD, 0xF5, 0xEC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CORAL     = RGBColor(0xF5, 0x63, 0x6E)
CORAL_DIM = RGBColor(0xF9, 0xA3, 0xA9)
YELLOW    = RGBColor(0xF7, 0xDF, 0x93)
BLUE      = RGBColor(0x2B, 0x36, 0x74)
INK       = RGBColor(0x26, 0x2B, 0x4A)
MUTED_L   = RGBColor(0x70, 0x76, 0x94)
MUTED_D   = RGBColor(0xA9, 0xAF, 0xCE)
FAINT_D   = RGBColor(0x6C, 0x73, 0x99)
LINE_L    = RGBColor(0xE7, 0xDF, 0xD6)
CHIP_L    = RGBColor(0xEF, 0xF2, 0xFA)

SERIF = "Bookman Old Style"
SANS  = "Calibri"

SW, SH = 13.333, 7.5
ML     = 0.62               # left margin
CW     = SW - 2 * ML        # content width

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

_page = {"n": 0}


# ------------------------------------------------------------------- primitives
def P(text, size=12, bold=False, color=None, font=SANS, space_after=5,
      line=1.15, align=None, italic=False, space_before=0):
    return dict(text=text, size=size, bold=bold, color=color, font=font,
                space_after=space_after, line=line, align=align,
                italic=italic, space_before=space_before)


def tx(container, l, t, w, h, paras, anchor=MSO_ANCHOR.TOP, default_color=INK):
    """Add a textbox (or fill an existing shape's text frame) with paragraphs."""
    if isinstance(container, tuple):
        slide = container[0]
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
    else:
        box = container
        tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spec["line"]
        p.space_after = Pt(spec["space_after"])
        p.space_before = Pt(spec["space_before"])
        if spec["align"]:
            p.alignment = spec["align"]
        r = p.add_run()
        r.text = spec["text"]
        f = r.font
        f.name = spec["font"]
        f.size = Pt(spec["size"])
        f.bold = spec["bold"]
        f.italic = spec["italic"]
        f.color.rgb = spec["color"] if spec["color"] is not None else default_color
    return box


def rect(slide, l, t, w, h, fill=None, radius=0.10, line=None, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = min(0.5, radius / min(w, h))
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    s.text_frame.margin_left = s.text_frame.margin_right = Inches(0.16)
    s.text_frame.margin_top = s.text_frame.margin_bottom = Inches(0.10)
    return s


def bar(slide, l, t, w, h, fill):
    return rect(slide, l, t, w, h, fill=fill, shape=MSO_SHAPE.RECTANGLE)


def circle(slide, l, t, d, fill):
    return rect(slide, l, t, d, d, fill=fill, shape=MSO_SHAPE.OVAL)


def numbered(slide, cl, ct, d, n, fill, fg=WHITE, size=15):
    c = circle(slide, cl, ct, d, fill)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(n)
    r.font.name = SERIF
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    return c


def slide_new(dark=False, numbered_page=True):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY if dark else CREAM
    s._dark = dark
    if numbered_page:
        _page["n"] += 1
        foot(s, dark, _page["n"])
    return s


def foot(s, dark, n):
    col = FAINT_D if dark else RGBColor(0xA8, 0xA0, 0x98)
    tx((s,), ML, 6.94, 6.0, 0.3,
       [P("Emotion Recognition System  \u2014  Convolutional Neural Network", 9,
          color=col, space_after=0)])
    tx((s,), SW - ML - 1.0, 6.94, 1.0, 0.3,
       [P(str(n), 9, color=col, space_after=0, align=PP_ALIGN.RIGHT)])


def head(s, title, sub=None):
    dark = getattr(s, "_dark", False)
    tx((s,), ML, 0.44, CW, 0.72,
       [P(title, 33, bold=True, color=WHITE if dark else INK, font=SERIF,
          space_after=0, line=1.0)])
    y = 1.15
    if sub:
        tx((s,), ML, 1.13, CW, 0.32,
           [P(sub, 12.5, color=YELLOW if dark else MUTED_L, space_after=0)])
        y = 1.50
    bar(s, ML, y, 0.44, 0.055, CORAL)
    return y + 0.34


# ================================================================== SLIDE 1
s = slide_new(dark=True, numbered_page=False)
circle(s, 9.9, -1.5, 4.6, RGBColor(0x3D, 0x2C, 0x46))
circle(s, 11.4, 5.2, 3.9, RGBColor(0x33, 0x38, 0x5A))

tx((s,), 1.0, 2.55, 10.8, 1.0,
   [P("EMOTION RECOGNITION SYSTEM", 40, bold=True, color=WHITE, font=SERIF,
      space_after=0, line=1.0)])
tx((s,), 1.0, 3.68, 10.6, 0.4,
   [P("Facial Expression Classification with a Convolutional Neural Network",
      17, color=CORAL, space_after=0)])
bar(s, 1.0, 4.42, 0.6, 0.055, CORAL)
tx((s,), 1.0, 4.78, 10.6, 0.35,
   [P("Final Year Project   |   CNN trained from scratch on FER-2013   |   "
      "57.09% accuracy on 7,178 unseen test images",
      11, color=MUTED_D, space_after=0)])

tx((s,), 1.0, 5.42, 4.6, 0.3,
   [P("P R E S E N T E D   B Y", 9, bold=True, color=YELLOW, space_after=0)])
tx((s,), 1.0, 5.80, 4.6, 1.2,
   [P("Aleeza Zubair  \u00b7  70147042", 12.5, color=WHITE, space_after=7),
    P("Sehar Mehmood  \u00b7  70147004", 12.5, color=WHITE, space_after=7),
    P("Aneeb Bari  \u00b7  70150957", 12.5, color=WHITE, space_after=0)])
tx((s,), 7.6, 5.42, 4.6, 0.3,
   [P("S U P E R V I S E D   B Y", 9, bold=True, color=YELLOW, space_after=0)])
tx((s,), 7.6, 5.80, 4.6, 0.4,
   [P("Amish Hassan", 12.5, color=WHITE, space_after=0)])


# ================================================================== SLIDE 2
s = slide_new()
y = head(s, "Introduction", "Why the automatic reading of facial expression is worth solving")

tx((s,), ML, y, 7.0, 2.4,
   [P("Facial expressions are one of the primary channels of non-verbal human "
      "communication. Automating their interpretation enables applications in "
      "human\u2013computer interaction, driver-attention monitoring, market research, "
      "mental-health screening and adaptive learning systems.", 13, line=1.35,
      space_after=13),
    P("The task is non-trivial. Expressions vary across individuals, cultures and "
      "contexts; they are often subtle or blended; and image conditions such as "
      "lighting, pose and occlusion vary widely.", 13, line=1.35, space_after=0)])

c = rect(s, 8.0, y - 0.05, 4.71, 2.35, fill=NAVY, radius=0.14)
tx(c, 0, 0, 0, 0,
   [P("The ceiling is not 100%", 15, bold=True, color=YELLOW, font=SERIF,
      space_after=10),
    P("Human annotators agree on FER-2013 labels only about 65% of the time. "
      "The ground truth is itself subjective, which places a practical ceiling "
      "on the accuracy any model can reach on this dataset.", 12, color=WHITE,
      line=1.3, space_after=0)],
   anchor=MSO_ANCHOR.MIDDLE)

chips = ["Human\u2013Computer Interaction", "Driver-Attention Monitoring",
         "Mental-Health Screening", "Adaptive Learning Systems"]
cw, gap = 2.90, 0.17
for i, t in enumerate(chips):
    x = ML + i * (cw + gap)
    ch = rect(s, x, 5.10, cw, 0.95, fill=WHITE, radius=0.10)
    bar(s, x, 5.10, cw, 0.055, CORAL)
    tx(ch, 0, 0, 0, 0,
       [P(t, 12, bold=True, color=INK, space_after=0, align=PP_ALIGN.CENTER)],
       anchor=MSO_ANCHOR.MIDDLE)

tx((s,), ML, 6.30, CW, 0.4,
   [P("Traditional approaches relied on hand-engineered features. A CNN learns the "
      "discriminative features directly from pixels \u2014 which is why it is the "
      "architecture used throughout this project.", 11.5, italic=True,
      color=MUTED_L, space_after=0)])


# ================================================================== SLIDE 3
s = slide_new(dark=True)
y = head(s, "Objectives")

objs = [
    "Build and train a CNN classifier for seven-class facial emotion recognition",
    "Evaluate the model rigorously using accuracy, precision, recall, F1-score and a confusion matrix",
    "Deploy the model in a real-time application driven by a live camera feed",
    "Analyse the failure modes of the system and identify concrete improvements",
]
rh, rg = 1.05, 0.20
for i, t in enumerate(objs):
    yy = y + 0.12 + i * (rh + rg)
    rect(s, ML, yy, CW, rh, fill=CARD_D, radius=0.12)
    numbered(s, ML + 0.34, yy + (rh - 0.62) / 2, 0.62, i + 1, CORAL)
    tx((s,), ML + 1.24, yy, CW - 1.6, rh,
       [P(t, 14.5, color=WHITE, space_after=0, line=1.25)],
       anchor=MSO_ANCHOR.MIDDLE)


# ================================================================== SLIDE 4
s = slide_new()
y = head(s, "Scope of the Project", "What the system classifies \u2014 and what was deliberately left out")

def scope_card(x, w, title, tcol, items, fill, fg, dot):
    c = rect(s, x, y, w, 4.35, fill=fill, radius=0.14)
    tx((s,), x + 0.38, y + 0.38, w - 0.76, 0.35,
       [P(title, 15, bold=True, color=tcol, font=SERIF, space_after=0)])
    for i, it in enumerate(items):
        yy = y + 1.00 + i * 0.78
        circle(s, x + 0.40, yy + 0.10, 0.13, dot)
        tx((s,), x + 0.72, yy - 0.02, w - 1.16, 0.7,
           [P(it, 12.5, color=fg, space_after=0, line=1.28)])

scope_card(ML, 5.90, "In Scope", CORAL, [
    "Single largest detected face per frame",
    "Seven fixed emotion categories",
    "48\u00d748 grayscale, frontal faces",
    "Real-time deployment in a browser",
], WHITE, INK, CORAL)

scope_card(ML + 6.20, 5.89, "Out of Scope", YELLOW, [
    "Multi-face tracking within a frame",
    "Emotion-intensity regression",
    "Profile views and heavy occlusion",
    "Remote / client-side camera capture",
], NAVY, WHITE, YELLOW)


# ================================================================== SLIDE 5
s = slide_new(dark=True)
y = head(s, "Background", "The benchmark, and why a convolutional network suits it")

c = rect(s, ML, y, 5.90, 3.05, fill=CARD_D, radius=0.14)
tx((s,), ML + 0.38, y + 0.34, 5.14, 0.32,
   [P("FER-2013", 15, bold=True, color=YELLOW, font=SERIF, space_after=0)])
tx((s,), ML + 0.38, y + 0.82, 5.14, 2.1,
   [P("Introduced for the ICML 2013 Challenges in Representation Learning. "
      "48\u00d748 grayscale face images collected by automated web search and "
      "labelled into seven emotion categories.", 12, color=WHITE, line=1.3,
      space_after=11),
    P("Its known difficulties \u2014 low resolution, label noise and heavy class "
      "imbalance \u2014 appear directly in the results of this project.", 12,
      color=MUTED_D, line=1.3, space_after=0)])

c = rect(s, ML + 6.20, y, 5.89, 3.05, fill=CARD_D, radius=0.14)
tx((s,), ML + 6.58, y + 0.34, 5.13, 0.32,
   [P("Why a CNN", 15, bold=True, color=YELLOW, font=SERIF, space_after=0)])
tx((s,), ML + 6.58, y + 0.82, 5.13, 2.1,
   [P("Early convolutional layers respond to edges and simple textures; deeper "
      "layers compose these into parts \u2014 eyes, mouth corners, brow furrows \u2014 "
      "that are discriminative for expression.", 12, color=WHITE, line=1.3,
      space_after=11),
    P("Weight sharing makes the representation translation-tolerant and cuts the "
      "parameter count far below a dense network of comparable capacity.", 12,
      color=MUTED_D, line=1.3, space_after=0)])

tx((s,), ML, y + 3.35, CW, 0.3,
   [P("W H E R E   T H I S   W O R K   S I T S", 9, bold=True, color=YELLOW,
      space_after=0)])
bench = [("\u224871%", "2013 competition winner", CARD_D2, WHITE),
         ("\u224865 \u00b1 5%", "Human performance", CARD_D2, WHITE),
         ("57.09%", "This project", CORAL, WHITE)]
bw = (CW - 2 * 0.22) / 3
for i, (big, lab, fill, fg) in enumerate(bench):
    x = ML + i * (bw + 0.22)
    rect(s, x, y + 3.72, bw, 1.02, fill=fill, radius=0.12)
    tx((s,), x, y + 3.86, bw, 0.4,
       [P(big, 20, bold=True, color=YELLOW if fill == CARD_D2 else WHITE,
          font=SERIF, space_after=0, align=PP_ALIGN.CENTER)])
    tx((s,), x, y + 4.32, bw, 0.3,
       [P(lab, 11, color=MUTED_D if fill == CARD_D2 else WHITE, space_after=0,
          align=PP_ALIGN.CENTER)])


# ================================================================== SLIDE 6
s = slide_new()
y = head(s, "Dataset \u2014 FER-2013", "35,887 labelled images across seven emotion categories")

stats = [("35,887", "Total images"), ("48\u00d748", "Pixels per image"),
         ("Grayscale", "Colour format"), ("7", "Emotion classes")]
sw_ = (CW - 3 * 0.20) / 4
for i, (big, lab) in enumerate(stats):
    x = ML + i * (sw_ + 0.20)
    rect(s, x, y, sw_, 1.10, fill=WHITE, radius=0.12)
    tx((s,), x, y + 0.20, sw_, 0.45,
       [P(big, 21, bold=True, color=CORAL, font=SERIF, space_after=0,
          align=PP_ALIGN.CENTER)])
    tx((s,), x, y + 0.72, sw_, 0.3,
       [P(lab, 10.5, color=MUTED_L, space_after=0, align=PP_ALIGN.CENTER)])

ty = y + 1.42
rows = [("Emotion", "Train images", "Test images", "Train share"),
        ("Angry", "3,995", "958", "13.9%"),
        ("Disgust", "436", "111", "1.5%"),
        ("Fear", "4,097", "1,024", "14.3%"),
        ("Happy", "7,215", "1,774", "25.1%"),
        ("Neutral", "4,965", "1,233", "17.3%"),
        ("Sad", "4,830", "1,247", "16.8%"),
        ("Surprise", "3,171", "831", "11.0%"),
        ("Total", "28,709", "7,178", "100%")]
tblw = 7.55
tbl = s.shapes.add_table(len(rows), 4, Inches(ML), Inches(ty),
                         Inches(tblw), Inches(3.55)).table
tbl.first_row = False
tbl.horz_banding = False
for w_, i in zip([2.35, 1.85, 1.75, 1.60], range(4)):
    tbl.columns[i].width = Inches(w_)
for ri, row in enumerate(rows):
    tbl.rows[ri].height = Inches(0.39)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = NAVY
        elif ri == len(rows) - 1:
            cell.fill.fore_color.rgb = CHIP_L
        else:
            cell.fill.fore_color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.13)
        cell.margin_top = cell.margin_bottom = 0
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = val
        r.font.name = SANS
        r.font.size = Pt(11.5)
        r.font.bold = (ri == 0 or ri == len(rows) - 1)
        r.font.color.rgb = WHITE if ri == 0 else INK

# split panel
px = ML + tblw + 0.30
pw = CW - tblw - 0.30
c = rect(s, px, ty, pw, 3.55, fill=NAVY, radius=0.14)
tx((s,), px + 0.34, ty + 0.32, pw - 0.68, 0.3,
   [P("How the data is split", 14, bold=True, color=YELLOW, font=SERIF,
      space_after=0)])

seg = [("Training", "22,967", 0.64, CORAL),
       ("Validation", "5,742", 0.16, YELLOW),
       ("Test", "7,178", 0.20, RGBColor(0x53, 0x5C, 0x94))]
bx, bw2 = px + 0.34, pw - 0.68
for lab, cnt, frac, col in seg:
    bar(s, bx, ty + 0.90, bw2 * frac, 0.34, col)
    bx += bw2 * frac
tx((s,), px + 0.34, ty + 1.38, pw - 0.68, 2.10,
   [P("Training  22,967  (64%)", 11.5, bold=True, color=CORAL, space_after=4),
    P("Validation  5,742  (16%)", 11.5, bold=True, color=YELLOW, space_after=4),
    P("Test  7,178  (20%)", 11.5, bold=True, color=MUTED_D, space_after=10),
    P("The 28,709 training images are split 80 / 20 into fit and validation "
      "partitions. The split is stratified with random_state = 42, so class "
      "proportions are preserved and validation metrics are not distorted by "
      "sampling.", 11, color=WHITE, line=1.28, space_after=0)])


# ================================================================== SLIDE 7
s = slide_new(dark=True)
y = head(s, "Class Imbalance", "The defining property of FER-2013 \u2014 and the source of the largest failure mode")

counts = [("Happy", 7215), ("Neutral", 4965), ("Sad", 4830), ("Fear", 4097),
          ("Angry", 3995), ("Surprise", 3171), ("Disgust", 436)]
maxc = 7215
bx, bwmax = ML + 1.30, 5.70
for i, (name, n) in enumerate(counts):
    yy = y + 0.10 + i * 0.55
    tx((s,), ML, yy - 0.01, 1.20, 0.3,
       [P(name, 12, bold=True, color=WHITE, space_after=0, align=PP_ALIGN.RIGHT)])
    col = CORAL if name == "Happy" else (YELLOW if name == "Disgust" else CARD_D2)
    bar(s, bx, yy, max(0.06, bwmax * n / maxc), 0.34, col)
    tx((s,), bx + max(0.06, bwmax * n / maxc) + 0.12, yy - 0.01, 1.1, 0.3,
       [P(f"{n:,}", 11.5, color=MUTED_D, space_after=0)])

c = rect(s, ML + 7.55, y, CW - 7.55, 3.95, fill=CARD_D, radius=0.14)
tx((s,), ML + 7.90, y + 0.36, CW - 8.25, 3.3,
   [P("16.5 : 1", 26, bold=True, color=CORAL, font=SERIF, space_after=4),
    P("Happy (7,215) against Disgust (436)", 11.5, color=MUTED_D, space_after=16),
    P("The network is trained with an unweighted cross-entropy loss. It can "
      "therefore reduce that loss more effectively by learning the frequent "
      "classes well and effectively ignoring the rare ones.", 12, color=WHITE,
      line=1.3, space_after=12),
    P("The results section shows that this is precisely what happened: Disgust "
      "recall falls to 0.171.", 12, bold=True, color=YELLOW, line=1.3,
      space_after=0)])


# ================================================================== SLIDE 8
s = slide_new()
y = head(s, "Data Preprocessing", "training/preprocessing.py \u2014 applied identically at training and inference time")

steps = [("Grayscale", "Each image is read in grayscale. Expression is carried by "
          "shape and shading, and this cuts input dimensionality threefold."),
         ("Resize", "Every image is resized to a uniform 48\u00d748 pixels."),
         ("Normalise", "Pixel values are scaled from [0, 255] to [0, 1], keeping "
          "gradients well-scaled during optimisation."),
         ("Reshape", "A channel dimension is added, giving each sample the shape "
          "(48, 48, 1)."),
         ("Cache", "Processed arrays are written to X_data.npy and y_labels.npy so "
          "training runs never repeat the decode step.")]
cw2 = (CW - 4 * 0.16) / 5
for i, (t, d) in enumerate(steps):
    x = ML + i * (cw2 + 0.16)
    rect(s, x, y, cw2, 2.75, fill=WHITE, radius=0.12)
    numbered(s, x + 0.28, y + 0.26, 0.52, i + 1, CORAL)
    tx((s,), x + 0.28, y + 0.94, cw2 - 0.56, 0.3,
       [P(t, 13.5, bold=True, color=INK, font=SERIF, space_after=0)])
    tx((s,), x + 0.28, y + 1.32, cw2 - 0.56, 1.3,
       [P(d, 11, color=MUTED_L, line=1.3, space_after=0)])
    if i < 4:
        tx((s,), x + cw2 + 0.005, y + 1.15, 0.15, 0.3,
           [P("\u203a", 20, bold=True, color=CORAL_DIM, space_after=0,
              align=PP_ALIGN.CENTER)])

wy = y + 3.00
c = rect(s, ML, wy, CW, 1.15, fill=NAVY, radius=0.12)
bar(s, ML, wy, 0.075, 1.15, YELLOW)
tx((s,), ML + 0.42, wy + 0.20, CW - 0.84, 0.8,
   [P("The label order is load-bearing", 13, bold=True, color=YELLOW,
      space_after=6),
    P("Labels are assigned by index into  [ angry, disgust, fear, happy, neutral, "
      "sad, surprise ].  Every component that maps a model output index back to a "
      "human-readable name must use exactly this order \u2014 a mismatch produces "
      "confident but wrong labels, silently.", 11.5, color=WHITE, line=1.3,
      space_after=0)])


# ================================================================== SLIDE 9
s = slide_new(dark=True)
y = head(s, "CNN Model Architecture", "training/model.py \u2014 355,847 trainable parameters")

layers = [("Input", "48 \u00d7 48 \u00d7 1  grayscale image", CARD_D, WHITE),
          ("Conv2D 32 @ 3\u00d73 + ReLU", "\u2192  MaxPooling 2\u00d72", CARD_D, WHITE),
          ("Conv2D 64 @ 3\u00d73 + ReLU", "\u2192  MaxPooling 2\u00d72", CARD_D, WHITE),
          ("Conv2D 128 @ 3\u00d73 + ReLU", "\u2192  MaxPooling 2\u00d72", CARD_D, WHITE),
          ("Flatten", "feature maps \u2192 1-D vector", CARD_D, WHITE),
          ("Dense 128 + ReLU", "\u2192  Dropout 0.5", CARD_D, WHITE),
          ("Dense 7 + Softmax", "probability over seven emotions", CORAL, WHITE)]
lh, lg = 0.62, 0.11
for i, (nm, det, fill, fg) in enumerate(layers):
    yy = y + i * (lh + lg)
    rect(s, ML, yy, 7.35, lh, fill=fill, radius=0.10)
    tx((s,), ML + 0.30, yy, 3.55, lh,
       [P(nm, 12.5, bold=True, color=fg, space_after=0)], anchor=MSO_ANCHOR.MIDDLE)
    tx((s,), ML + 3.95, yy, 3.10, lh,
       [P(det, 11, color=WHITE if fill == CORAL else MUTED_D, space_after=0)],
       anchor=MSO_ANCHOR.MIDDLE)

ax = ML + 7.75
notes = [("355,847", "trainable parameters", CORAL),
         ("32 \u2192 64 \u2192 128", "Filter counts double at each stage while pooling "
          "halves the spatial resolution \u2014 trading spatial detail for "
          "representational depth.", None),
         ("One regulariser", "A single Dropout(0.5) before the classifier is the "
          "only explicit regularisation in the network. This is revisited under "
          "Limitations.", None)]
c = rect(s, ax, y, CW - 7.75, 1.30, fill=CARD_D2, radius=0.12)
tx((s,), ax + 0.30, y + 0.26, CW - 8.35, 0.9,
   [P("355,847", 25, bold=True, color=YELLOW, font=SERIF, space_after=2),
    P("trainable parameters", 11, color=MUTED_D, space_after=0)])
c = rect(s, ax, y + 1.44, CW - 7.75, 1.55, fill=CARD_D, radius=0.12)
tx((s,), ax + 0.30, y + 1.68, CW - 8.35, 1.1,
   [P("32 \u2192 64 \u2192 128", 13, bold=True, color=CORAL, space_after=6),
    P("Filter counts double at each stage while pooling halves the spatial "
      "resolution \u2014 a standard design that trades spatial detail for "
      "representational depth.", 11, color=WHITE, line=1.28, space_after=0)])
c = rect(s, ax, y + 3.13, CW - 7.75, 1.55, fill=CARD_D, radius=0.12)
tx((s,), ax + 0.30, y + 3.37, CW - 8.35, 1.1,
   [P("A single regulariser", 13, bold=True, color=CORAL, space_after=6),
    P("Dropout(0.5) before the classifier is the only explicit regularisation in "
      "the network \u2014 a point returned to under Limitations.", 11, color=WHITE,
      line=1.28, space_after=0)])


# ================================================================== SLIDE 10
s = slide_new()
y = head(s, "Training Configuration", "training/train.py")

cfg = [("Optimiser", "Adam, default learning rate 0.001"),
       ("Loss function", "Categorical cross-entropy"),
       ("Batch size", "64 images per step"),
       ("Maximum epochs", "50"),
       ("Validation split", "20% of the training set, stratified, random_state = 42"),
       ("Early stopping", "Monitor val_loss, patience 10, restore best weights"),
       ("Learning-rate reduction", "Factor 0.5, patience 5, minimum 1e-5"),
       ("Checkpointing", "Save the best model by val_accuracy")]
rh2 = 0.50
for i, (k, v) in enumerate(cfg):
    yy = y + i * (rh2 + 0.10)
    rect(s, ML, yy, 7.90, rh2, fill=WHITE, radius=0.08)
    tx((s,), ML + 0.28, yy, 2.90, rh2,
       [P(k, 12, bold=True, color=INK, space_after=0)], anchor=MSO_ANCHOR.MIDDLE)
    tx((s,), ML + 3.30, yy, 4.35, rh2,
       [P(v, 11.5, color=MUTED_L, space_after=0)], anchor=MSO_ANCHOR.MIDDLE)

nx = ML + 8.30
c = rect(s, nx, y, CW - 8.30, 2.20, fill=NAVY, radius=0.14)
tx((s,), nx + 0.32, y + 0.30, CW - 8.94, 1.7,
   [P("Why stratify?", 14, bold=True, color=YELLOW, font=SERIF, space_after=8),
    P("The stratified split preserves the class proportions of the full dataset in "
      "both the training and validation partitions, so validation metrics are not "
      "distorted by sampling.", 11.5, color=WHITE, line=1.3, space_after=0)])
c = rect(s, nx, y + 2.40, CW - 8.30, 2.20, fill=WHITE, radius=0.14)
tx((s,), nx + 0.32, y + 2.70, CW - 8.94, 1.7,
   [P("Why early stopping?", 14, bold=True, color=CORAL, font=SERIF, space_after=8),
    P("Training was permitted to run for 50 epochs but was halted automatically at "
      "epoch 19, with the best weights restored \u2014 the network had already begun "
      "to overfit.", 11.5, color=MUTED_L, line=1.3, space_after=0)])


# ================================================================== SLIDE 11
s = slide_new()
y = head(s, "Training Behaviour", "What the learning curves reveal")

facts = [("Epoch 19 / 50", "Early stopping fired; best weights restored"),
         ("\u2248 72%", "Final training accuracy"),
         ("\u2248 56%", "Validation accuracy \u2014 plateaued and flat"),
         ("Epoch 8", "Validation loss bottoms at \u22481.20, then rises to \u22481.29")]
fw = (CW - 3 * 0.18) / 4
for i, (big, lab) in enumerate(facts):
    x = ML + i * (fw + 0.18)
    rect(s, x, y, fw, 0.96, fill=WHITE, radius=0.12)
    bar(s, x, y, fw, 0.05, CORAL)
    tx((s,), x + 0.22, y + 0.16, fw - 0.44, 0.30,
       [P(big, 14.5, bold=True, color=INK, font=SERIF, space_after=0)])
    tx((s,), x + 0.22, y + 0.52, fw - 0.44, 0.42,
       [P(lab, 10, color=MUTED_L, line=1.18, space_after=0)])

ih = 3.15
iw = ih * 3.062
ix = (SW - iw) / 2
iy = y + 1.06
if os.path.exists(HIST):
    s.shapes.add_picture(HIST, Inches(ix), Inches(iy), Inches(iw), Inches(ih))

cy = 6.16
c = rect(s, ML, cy, CW, 0.68, fill=NAVY, radius=0.10)
bar(s, ML, cy, 0.07, 0.68, YELLOW)
tx((s,), ML + 0.40, cy, CW - 0.80, 0.68,
   [P("A rising validation loss alongside a falling training loss is the textbook "
      "signature of overfitting: beyond roughly epoch 8 the network is memorising "
      "training examples rather than learning generalisable features.", 11.5,
      color=WHITE, line=1.25, space_after=0)], anchor=MSO_ANCHOR.MIDDLE)


# ================================================================== SLIDE 12
s = slide_new(dark=True)
y = head(s, "Results \u2014 Overall Performance",
         "Evaluated on the complete 7,178-image test set, disjoint from training and validation")

c = rect(s, ML, y, 4.60, 2.55, fill=CORAL, radius=0.16)
tx((s,), ML, y + 0.42, 4.60, 1.0,
   [P("57.09%", 52, bold=True, color=WHITE, font=SERIF, space_after=0,
      align=PP_ALIGN.CENTER)])
tx((s,), ML, y + 1.62, 4.60, 0.6,
   [P("Test accuracy across seven classes\non 7,178 unseen images", 12.5,
      color=WHITE, space_after=0, align=PP_ALIGN.CENTER, line=1.3)])

for i, (big, lab) in enumerate([("0.514", "Macro-averaged F1"),
                                ("0.566", "Weighted-averaged F1")]):
    x = ML + 4.90 + i * 2.30
    rect(s, x, y, 2.10, 2.55, fill=CARD_D, radius=0.16)
    tx((s,), x, y + 0.82, 2.10, 0.5,
       [P(big, 26, bold=True, color=YELLOW, font=SERIF, space_after=0,
          align=PP_ALIGN.CENTER)])
    tx((s,), x, y + 1.44, 2.10, 0.4,
       [P(lab, 11, color=MUTED_D, space_after=0, align=PP_ALIGN.CENTER,
          line=1.25)])

bxx = ML + 9.55
c = rect(s, bxx, y, CW - 9.55, 2.55, fill=CARD_D2, radius=0.16)
tx((s,), bxx + 0.32, y + 0.30, CW - 10.19, 2.0,
   [P("Against baselines", 13.5, bold=True, color=YELLOW, font=SERIF,
      space_after=10),
    P("Random guess (1 / 7)          14.29%", 11.5, color=MUTED_D, space_after=7),
    P("Always predict \u201chappy\u201d       24.71%", 11.5, color=MUTED_D, space_after=7),
    P("This model                          57.09%", 11.5, bold=True, color=CORAL,
      space_after=0)])

byy = y + 2.90
labels = [("Random baseline (1 / 7)", 14.29, CARD_D2, MUTED_D),
          ("Majority class \u2014 always \u201chappy\u201d", 24.71, CARD_D2, MUTED_D),
          ("Trained CNN", 57.09, CORAL, WHITE)]
for i, (lab, val, col, fg) in enumerate(labels):
    yy = byy + i * 0.60
    tx((s,), ML, yy - 0.01, 3.55, 0.3,
       [P(lab, 11.5, color=WHITE, space_after=0, align=PP_ALIGN.RIGHT)])
    bar(s, ML + 3.75, yy, 7.20 * val / 60.0, 0.36, col)
    tx((s,), ML + 3.75 + 7.20 * val / 60.0 + 0.14, yy - 0.01, 1.2, 0.3,
       [P(f"{val:.2f}%", 11.5, bold=True, color=fg if col == CORAL else MUTED_D,
          space_after=0)])

tx((s,), ML, byy + 1.92, CW, 0.35,
   [P("The model performs well above both baselines, confirming that it has "
      "learned genuine discriminative structure rather than exploiting the class "
      "prior.", 11.5, italic=True, color=MUTED_D, space_after=0)])


# ================================================================== SLIDE 13
s = slide_new()
y = head(s, "Per-Class Performance", "Accuracy alone conceals how uneven the model is")

prows = [("Emotion", "Precision", "Recall", "F1", "Support"),
         ("Angry", "0.494", "0.478", "0.486", "958"),
         ("Disgust", "0.792", "0.171", "0.281", "111"),
         ("Fear", "0.422", "0.294", "0.346", "1,024"),
         ("Happy", "0.767", "0.796", "0.781", "1,774"),
         ("Neutral", "0.496", "0.606", "0.545", "1,233"),
         ("Sad", "0.423", "0.477", "0.448", "1,247"),
         ("Surprise", "0.747", "0.681", "0.712", "831")]
strong = {"Happy", "Surprise"}
weak = {"Fear", "Sad", "Disgust"}
tw = 7.10
tbl = s.shapes.add_table(len(prows), 5, Inches(ML), Inches(y),
                         Inches(tw), Inches(3.60)).table
tbl.first_row = False
tbl.horz_banding = False
for w_, i in zip([1.72, 1.40, 1.32, 1.30, 1.36], range(5)):
    tbl.columns[i].width = Inches(w_)
for ri, row in enumerate(prows):
    tbl.rows[ri].height = Inches(0.44)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.12)
        cell.margin_top = cell.margin_bottom = 0
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = val
        r.font.name = SANS
        r.font.size = Pt(11.5)
        if ri == 0:
            r.font.bold = True
            r.font.color.rgb = WHITE
        else:
            nm = row[0]
            r.font.bold = nm in strong and ci == 3
            if nm in strong and ci == 3:
                r.font.color.rgb = RGBColor(0x1B, 0x7F, 0x5A)
            elif nm in weak and ci == 3:
                r.font.color.rgb = CORAL
            else:
                r.font.color.rgb = INK

ix = ML + tw + 0.30
iwid = CW - tw - 0.30
cards = [("Recognised reliably", CORAL,
          "Happy (F1 0.781) and Surprise (F1 0.712). Both involve large, "
          "unambiguous geometric deformation \u2014 a raised mouth curve, lifted "
          "brows \u2014 that survives downsampling to 48\u00d748.", WHITE, INK, 1.42),
         ("Recognised poorly", CORAL,
          "Fear (F1 0.346) and Sad (F1 0.448) are subtle, low-contrast "
          "expressions whose distinguishing cues are fine-grained and largely "
          "destroyed at this resolution.", WHITE, INK, 1.30),
         ("Disgust \u2014 the revealing case", YELLOW,
          "Precision 0.792 but recall 0.171. When the model predicts disgust it is "
          "usually right, but it makes that prediction only 24 times in 7,178 "
          "images, catching 19 of 111 real cases. A direct consequence of the 1.5% "
          "training share \u2014 and invisible to accuracy alone.", NAVY, WHITE, 1.90)]
cyy = y
for title_, tcol, body, fill, fg, h_ in cards:
    rect(s, ix, cyy, iwid, h_, fill=fill, radius=0.12)
    tx((s,), ix + 0.26, cyy + 0.18, iwid - 0.52, 0.28,
       [P(title_, 12.5, bold=True, color=tcol, space_after=0)])
    tx((s,), ix + 0.26, cyy + 0.52, iwid - 0.52, h_ - 0.68,
       [P(body, 10.5, color=fg, line=1.26, space_after=0)])
    cyy += h_ + 0.16


# ================================================================== SLIDE 14
s = slide_new()
y = head(s, "Confusion Analysis", "The errors are systematic, not arbitrary")

ih2 = 4.35
iw2 = ih2 * 1.204
if os.path.exists(CONF):
    s.shapes.add_picture(CONF, Inches(ML), Inches(y), Inches(iw2), Inches(ih2))

tx2 = ML + iw2 + 0.40
tw2 = SW - ML - tx2
crows = [("True \u2192 Predicted", "Count", "% of class"),
         ("Disgust \u2192 Angry", "40", "36.0%"),
         ("Fear \u2192 Sad", "279", "27.2%"),
         ("Sad \u2192 Neutral", "279", "22.4%"),
         ("Neutral \u2192 Sad", "234", "19.0%"),
         ("Angry \u2192 Sad", "174", "18.2%"),
         ("Surprise \u2192 Fear", "106", "12.8%")]
tbl = s.shapes.add_table(len(crows), 3, Inches(tx2), Inches(y),
                         Inches(tw2), Inches(2.60)).table
tbl.first_row = False
tbl.horz_banding = False
for w_, i in zip([2.65, 1.05, 1.35], range(3)):
    tbl.columns[i].width = Inches(w_)
for ri, row in enumerate(crows):
    tbl.rows[ri].height = Inches(0.37)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.11)
        cell.margin_top = cell.margin_bottom = 0
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = val
        r.font.name = SANS
        r.font.size = Pt(11)
        r.font.bold = (ri == 0)
        r.font.color.rgb = WHITE if ri == 0 else INK

ny = y + 2.80
c = rect(s, tx2, ny, tw2, 1.55, fill=NAVY, radius=0.12)
tx((s,), tx2 + 0.28, ny + 0.24, tw2 - 0.56, 1.1,
   [P("Every dominant confusion pairs emotions that are genuinely similar in "
      "facial geometry and affective valence.", 11.5, bold=True, color=YELLOW,
      line=1.28, space_after=8),
    P("Disgust and anger share a lowered brow and raised upper lip. Fear and "
      "surprise share widened eyes and raised brows, differing mainly in mouth "
      "tension \u2014 a cue largely lost at 48\u00d748.", 10.5, color=WHITE, line=1.28,
      space_after=0)])

c = rect(s, tx2, ny + 1.70, tw2, 0.85, fill=WHITE, radius=0.12)
tx((s,), tx2 + 0.28, ny + 1.70, tw2 - 0.56, 0.85,
   [P("Sad, neutral and angry form a mutually confusable cluster of low-arousal "
      "expressions. The model\u2019s errors mirror the ambiguities human annotators "
      "also face.", 10.5, color=MUTED_L, line=1.28, space_after=0)],
   anchor=MSO_ANCHOR.MIDDLE)


# ================================================================== SLIDE 15
s = slide_new(dark=True)
y = head(s, "System Implementation", "Flask backend, Haar Cascade face localisation, CNN inference")

pipe = [("Face localisation",
         "A Haar Cascade locates faces in the grayscale frame. Where several are "
         "found, the largest bounding box is taken as the subject."),
        ("Crop",
         "The face region is cropped. Essential \u2014 the network was trained on "
         "tight crops, so a full scene shifts the input off-distribution."),
        ("Normalise",
         "The crop is resized to 48\u00d748 and scaled to [0, 1], matching training "
         "preprocessing exactly."),
        ("Predict",
         "Softmax yields a probability per class; the argmax is reported with its "
         "confidence score.")]
pw2 = (CW - 3 * 0.18) / 4
for i, (t, d) in enumerate(pipe):
    x = ML + i * (pw2 + 0.18)
    rect(s, x, y, pw2, 2.30, fill=CARD_D, radius=0.12)
    numbered(s, x + 0.26, y + 0.24, 0.50, i + 1, CORAL)
    tx((s,), x + 0.26, y + 0.88, pw2 - 0.52, 0.3,
       [P(t, 13, bold=True, color=WHITE, font=SERIF, space_after=0)])
    tx((s,), x + 0.26, y + 1.24, pw2 - 0.52, 0.98,
       [P(d, 10.5, color=MUTED_D, line=1.26, space_after=0)])
    if i < 3:
        tx((s,), x + pw2 + 0.01, y + 0.82, 0.16, 0.3,
           [P("\u203a", 20, bold=True, color=CORAL, space_after=0,
              align=PP_ALIGN.CENTER)])

ey = y + 2.52
erows = [("Endpoint", "Method", "Purpose"),
         ("/", "GET", "Web user interface"),
         ("/video_feed", "GET", "MJPEG stream of annotated frames"),
         ("/status", "GET", "Latest prediction and all seven confidence scores"),
         ("/camera/start, /camera/stop", "POST", "Acquire and release the webcam"),
         ("/api/health", "GET", "Service and model status")]
etw = 7.55
tbl = s.shapes.add_table(len(erows), 3, Inches(ML), Inches(ey),
                         Inches(etw), Inches(2.20)).table
tbl.first_row = False
tbl.horz_banding = False
for w_, i in zip([2.85, 1.05, 3.65], range(3)):
    tbl.columns[i].width = Inches(w_)
for ri, row in enumerate(erows):
    tbl.rows[ri].height = Inches(0.36)
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CORAL if ri == 0 else CARD_D
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.12)
        cell.margin_top = cell.margin_bottom = 0
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = val
        r.font.name = "Consolas" if (ci == 0 and ri > 0) else SANS
        r.font.size = Pt(10.5)
        r.font.bold = (ri == 0)
        r.font.color.rgb = WHITE

px2 = ML + etw + 0.30
c = rect(s, px2, ey, CW - etw - 0.30, 2.20, fill=CARD_D2, radius=0.12)
tx((s,), px2 + 0.32, ey + 0.28, CW - etw - 0.94, 1.7,
   [P("\u2248 10 fps", 22, bold=True, color=YELLOW, font=SERIF, space_after=6),
    P("Running the CNN on every captured frame stalls the stream. Inference "
      "therefore runs on every fifth frame and reuses the most recent prediction "
      "in between; face detection and box drawing still run every frame, so the "
      "overlay tracks smoothly. Measured at 84 frames over an 8-second capture.",
      10.5, color=WHITE, line=1.26, space_after=0)])


# ================================================================== SLIDE 16
s = slide_new()
y = head(s, "Real-Time Camera Detection", "The trained model running live in the browser")

pts = [("Server-side capture",
        "Flask serves a single-page application; OpenCV captures webcam frames on "
        "the machine running the server."),
       ("MJPEG streaming",
        "Annotated frames are JPEG-encoded and pushed as a multipart/x-mixed-replace "
        "stream, rendered by a plain <img> element \u2014 no client-side video or ML code "
        "is required."),
       ("Live confidence readout",
        "The browser polls /status four times per second to update the emotion label "
        "and the confidence bars for all seven classes.")]
for i, (t, d) in enumerate(pts):
    yy = y + i * 1.32
    rect(s, ML, yy, 5.85, 1.15, fill=WHITE, radius=0.12)
    bar(s, ML, yy, 0.06, 1.15, CORAL)
    tx((s,), ML + 0.34, yy + 0.20, 5.20, 0.28,
       [P(t, 13, bold=True, color=INK, font=SERIF, space_after=0)])
    tx((s,), ML + 0.34, yy + 0.54, 5.20, 0.55,
       [P(d, 10.5, color=MUTED_L, line=1.26, space_after=0)])

phx = ML + 6.20
ph = rect(s, phx, y, CW - 6.20, 4.30, fill=WHITE, radius=0.14,
          line=RGBColor(0xE0, 0xD6, 0xCB), line_w=1.25)
tx((s,), phx + 0.4, y + 1.75, CW - 7.0, 0.9,
   [P("[ Screenshot placeholder ]", 13, bold=True, color=RGBColor(0xC0, 0xB6, 0xAC),
      space_after=6, align=PP_ALIGN.CENTER),
    P("Paste your live-detection screenshot here \u2014 the annotated webcam frame "
      "with the bounding box and emotion label.", 10.5,
      color=RGBColor(0xC0, 0xB6, 0xAC), align=PP_ALIGN.CENTER, line=1.25,
      space_after=0)])


# ================================================================== SLIDE 17
s = slide_new(dark=True)
y = head(s, "Discussion and Limitations", "The two principal weaknesses, both identifiable in the design")

c = rect(s, ML, y, 5.90, 2.75, fill=CARD_D, radius=0.14)
tx((s,), ML + 0.34, y + 0.28, 5.22, 2.3,
   [P("1.  Overfitting", 15, bold=True, color=CORAL, font=SERIF, space_after=10),
    P("No data augmentation.  The network sees each of the 28,709 images in an "
      "identical form every epoch. Random flips, rotations, zooms and brightness "
      "shifts would multiply the effective dataset size at zero labelling cost.",
      11, color=WHITE, line=1.28, space_after=8),
    P("Minimal regularisation.  A single dropout layer before the classifier is "
      "the only explicit regulariser. Batch normalisation and weight decay would "
      "further constrain the network.", 11, color=MUTED_D, line=1.28,
      space_after=0)])

c = rect(s, ML + 6.20, y, 5.89, 2.75, fill=CARD_D, radius=0.14)
tx((s,), ML + 6.54, y + 0.28, 5.21, 2.3,
   [P("2.  Class Imbalance", 15, bold=True, color=CORAL, font=SERIF,
      space_after=10),
    P("Training with an unweighted loss allows the network to abandon the rare "
      "classes, as the disgust recall of 0.171 demonstrates.", 11, color=WHITE,
      line=1.28, space_after=8),
    P("Passing class_weight to model.fit() \u2014 inversely proportional to class "
      "frequency \u2014 would penalise errors on rare classes proportionally more, at "
      "some cost to majority-class accuracy. Oversampling and focal loss are "
      "alternative remedies.", 11, color=MUTED_D, line=1.28, space_after=0)])

oy = y + 3.05
tx((s,), ML, oy, CW, 0.3,
   [P("O T H E R   C O N S T R A I N T S", 9, bold=True, color=YELLOW,
      space_after=0)])
cons = [("Frontal faces only", "The Haar Cascade fails on profile views and heavy "
         "occlusion \u2014 no face means no prediction."),
        ("Server-side camera", "Capture happens on the machine running the server, "
         "which suits a single-machine demonstration."),
        ("Lighting sensitivity", "Training data offers limited illumination "
         "diversity; accuracy degrades in low light."),
        ("Single-face assumption", "Only the largest detected face is classified in "
         "each frame.")]
ow = (CW - 3 * 0.18) / 4
for i, (t, d) in enumerate(cons):
    x = ML + i * (ow + 0.18)
    rect(s, x, oy + 0.38, ow, 1.15, fill=CARD_D2, radius=0.10)
    tx((s,), x + 0.24, oy + 0.56, ow - 0.48, 0.28,
       [P(t, 11.5, bold=True, color=WHITE, space_after=0)])
    tx((s,), x + 0.24, oy + 0.86, ow - 0.48, 0.6,
       [P(d, 10, color=MUTED_D, line=1.22, space_after=0)])


# ================================================================== SLIDE 18
s = slide_new()
y = head(s, "A Defect Found and Corrected",
         "Label-order consistency \u2014 and why end-to-end verification mattered")

c = rect(s, ML, y, 7.40, 2.05, fill=NAVY, radius=0.14)
tx((s,), ML + 0.34, y + 0.26, 6.72, 0.3,
   [P("The mismatch", 13.5, bold=True, color=YELLOW, font=SERIF, space_after=0)])
tx((s,), ML + 0.34, y + 0.66, 6.72, 1.3,
   [P("API returned:      [ Angry, Disgust, Fear, Happy, Sad, Surprise, Neutral ]",
      11, color=CORAL, space_after=6),
    P("Model trained on:  [ Angry, Disgust, Fear, Happy, Neutral, Sad, Surprise ]",
      11, color=YELLOW, space_after=10),
    P("Indices 4, 5 and 6 were permuted, so every prediction falling into those "
      "three classes was reported under the wrong name.", 11.5, color=WHITE,
      line=1.28, space_after=0)])

qx = ML + 7.70
c = rect(s, qx, y, CW - 7.70, 2.05, fill=CORAL, radius=0.14)
tx((s,), qx, y + 0.30, CW - 7.70, 1.5,
   [P("50.0%   vs   27.6%", 20, bold=True, color=WHITE, font=SERIF,
      space_after=8, align=PP_ALIGN.CENTER),
    P("Accuracy on the same held-out images under the correct ordering versus the "
      "incorrect one. The fix nearly doubled measured accuracy.", 11, color=WHITE,
      line=1.26, space_after=0, align=PP_ALIGN.CENTER)])

wy2 = y + 2.35
c = rect(s, ML, wy2, CW, 0.85, fill=WHITE, radius=0.12)
bar(s, ML, wy2, 0.06, 0.85, YELLOW)
tx((s,), ML + 0.36, wy2, CW - 0.72, 0.85,
   [P("The defect was silent. The API returned well-formed responses with high "
      "confidence values, and nothing in the output indicated an error.", 12,
      bold=True, color=INK, space_after=0)], anchor=MSO_ANCHOR.MIDDLE)

ly = wy2 + 1.05
lessons = [("Lesson 1",
            "The class ordering is a contract between training and inference. It "
            "should be defined once and imported everywhere, never restated per file."),
           ("Lesson 2",
            "A defect of this kind is invisible to unit-level checks. It surfaces "
            "only when predictions are compared against known labels end to end.")]
lw = (CW - 0.24) / 2
for i, (t, d) in enumerate(lessons):
    x = ML + i * (lw + 0.24)
    rect(s, x, ly, lw, 1.30, fill=WHITE, radius=0.12)
    tx((s,), x + 0.32, ly + 0.24, lw - 0.64, 0.3,
       [P(t, 13, bold=True, color=CORAL, font=SERIF, space_after=0)])
    tx((s,), x + 0.32, ly + 0.62, lw - 0.64, 0.6,
       [P(d, 11.5, color=MUTED_L, line=1.28, space_after=0)])


# ================================================================== SLIDE 19
s = slide_new(dark=True)
y = head(s, "Future Work", "Ordered by expected benefit relative to effort")

fut = [("Data augmentation",
        "Random flips, rotations, zooms and brightness jitter. Directly addresses "
        "the primary limitation; expected to add several points of accuracy."),
       ("Class weighting",
        "Recover usable recall on disgust and fear by penalising rare-class errors "
        "proportionally more."),
       ("Batch normalisation",
        "Deeper regularisation for faster, more stable convergence and reduced "
        "overfitting."),
       ("Transfer learning",
        "Fine-tune a VGG-16 or ResNet-50 backbone pretrained on ImageNet. Published "
        "results reach \u224870% on FER-2013."),
       ("Better face detection",
        "Replace the Haar Cascade with an MTCNN or DNN-based detector for "
        "robustness to pose and lighting."),
       ("Temporal smoothing",
        "Average predictions over a short window of frames to suppress the flicker "
        "of per-frame classification.")]
gw = (CW - 2 * 0.20) / 3
gh = 2.12
for i, (t, d) in enumerate(fut):
    col, row = i % 3, i // 3
    x = ML + col * (gw + 0.20)
    yy = y + row * (gh + 0.22)
    rect(s, x, yy, gw, gh, fill=CARD_D, radius=0.12)
    numbered(s, x + 0.28, yy + 0.26, 0.50, i + 1,
             CORAL if i < 2 else CARD_D2, WHITE if i < 2 else YELLOW)
    tx((s,), x + 0.28, yy + 0.90, gw - 0.56, 0.3,
       [P(t, 13, bold=True, color=WHITE, font=SERIF, space_after=0)])
    tx((s,), x + 0.28, yy + 1.24, gw - 0.56, 0.78,
       [P(d, 10.5, color=MUTED_D, line=1.24, space_after=0)])


# ================================================================== SLIDE 20
s = slide_new()
y = head(s, "Application Extension \u2014 MoodFlick",
         "A demonstration of the trained classifier used as a component in a larger product")

pts2 = [("Same model, new surface",
         "The identical emotion_model.h5 classifier drives MoodFlick; nothing about "
         "the network or its training changes."),
        ("Emotion \u2192 recommendation",
         "A scanned expression is mapped to a curated movie shelf, with the "
         "prediction confidence carried through into the ranking."),
        ("Manual refinement",
         "Eight mood tags, fourteen genres and an actor/title search let users "
         "refine results beyond the camera scan.")]
for i, (t, d) in enumerate(pts2):
    yy = y + i * 1.20
    rect(s, ML, yy, 5.85, 1.05, fill=WHITE, radius=0.12)
    bar(s, ML, yy, 0.06, 1.05, CORAL)
    tx((s,), ML + 0.34, yy + 0.18, 5.20, 0.28,
       [P(t, 12.5, bold=True, color=INK, font=SERIF, space_after=0)])
    tx((s,), ML + 0.34, yy + 0.50, 5.20, 0.5,
       [P(d, 10.5, color=MUTED_L, line=1.24, space_after=0)])

c = rect(s, ML, y + 3.65, 5.85, 0.85, fill=NAVY, radius=0.12)
tx((s,), ML + 0.34, y + 3.65, 5.20, 0.85,
   [P("MoodFlick is presented as an application of the model, not as the "
      "contribution of this project. The classifier itself remains the work.",
      11, italic=True, color=YELLOW, line=1.26, space_after=0)],
   anchor=MSO_ANCHOR.MIDDLE)

phx2 = ML + 6.20
rect(s, phx2, y, CW - 6.20, 4.50, fill=WHITE, radius=0.14,
     line=RGBColor(0xE0, 0xD6, 0xCB), line_w=1.25)
tx((s,), phx2 + 0.4, y + 1.85, CW - 7.0, 0.9,
   [P("[ Screenshot placeholder ]", 13, bold=True, color=RGBColor(0xC0, 0xB6, 0xAC),
      space_after=6, align=PP_ALIGN.CENTER),
    P("Paste a MoodFlick screenshot here \u2014 the detected-emotion panel and the "
      "recommended shelf.", 10.5, color=RGBColor(0xC0, 0xB6, 0xAC),
      align=PP_ALIGN.CENTER, line=1.25, space_after=0)])


# ================================================================== SLIDE 21
s = slide_new(dark=True)
circle(s, -1.6, 4.4, 4.2, RGBColor(0x3D, 0x2C, 0x46))
y = head(s, "Conclusion")

tx((s,), ML, y + 0.10, 11.0, 2.4,
   [P("A complete facial emotion recognition system was designed, trained, "
      "evaluated and deployed. The CNN attains 57.09% accuracy across seven "
      "classes on 7,178 unseen test images \u2014 substantially exceeding the 14.29% "
      "random and 24.71% majority-class baselines, and comparable to typical "
      "results for a network of this size trained from scratch on FER-2013 "
      "without augmentation.", 14, color=WHITE, line=1.4, space_after=14),
    P("Beyond the headline figure, the evaluation established where the model "
      "succeeds and fails: high-signal expressions are recognised reliably, subtle "
      "and under-represented ones are not, and the training curves locate the "
      "onset of overfitting at approximately epoch 8 \u2014 identifying data "
      "augmentation and class weighting as the concrete next steps.", 13,
      color=MUTED_D, line=1.4, space_after=0)])

chips2 = [("57.09%", "test accuracy"), ("355,847", "parameters"),
          ("\u2248 10 fps", "live inference"), ("0.514", "macro F1")]
cw3 = (CW - 3 * 0.20) / 4
for i, (big, lab) in enumerate(chips2):
    x = ML + i * (cw3 + 0.20)
    rect(s, x, 4.55, cw3, 1.00, fill=CARD_D, radius=0.12)
    tx((s,), x, 4.72, cw3, 0.35,
       [P(big, 18, bold=True, color=YELLOW, font=SERIF, space_after=0,
          align=PP_ALIGN.CENTER)])
    tx((s,), x, 5.14, cw3, 0.3,
       [P(lab, 10.5, color=MUTED_D, space_after=0, align=PP_ALIGN.CENTER)])

tx((s,), ML, 5.90, 8.0, 0.6,
   [P("Thank You", 30, bold=True, color=YELLOW, font=SERIF, space_after=0)])


prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
